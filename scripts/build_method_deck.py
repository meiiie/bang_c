#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build the Neko Core "Thuyet minh phuong phap" deck from the storyboard.

Source of truth for content: docs/method-deck-storyboard-vi.md (which distils
docs/method-writeup-vi.md). Regenerable + version-controlled so the numbers stay
in sync. Light "engineering paper" theme. Fonts: Segoe UI (installed on every
Windows, correct Vietnamese diacritics) + Consolas for code/diagram.

Run:  python scripts/build_method_deck.py
Out:  docs/Neko-Core-Thuyet-minh-phuong-phap.pptx  (14 slides, with speaker notes)
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BANNER = os.path.join(ROOT, "docs", "assets", "neko-core-banner.png")
OUT = os.path.join(ROOT, "docs", "Neko-Core-Thuyet-minh-phuong-phap.pptx")

# ---- palette (light "engineering paper") -----------------------------------
BG      = RGBColor.from_string("F7F8FA")
INK_BG  = RGBColor.from_string("0B1220")
TEXT    = RGBColor.from_string("0F172A")
SUBTLE  = RGBColor.from_string("475569")
TEAL    = RGBColor.from_string("14B8A6")
INDIGO  = RGBColor.from_string("6366F1")
GREEN   = RGBColor.from_string("16A34A")
RED     = RGBColor.from_string("DC2626")
AMBER   = RGBColor.from_string("D97706")
WHITE   = RGBColor.from_string("FFFFFF")
BORDER  = RGBColor.from_string("E2E8F0")
SURF    = RGBColor.from_string("EEF2F6")
NAVYTX  = RGBColor.from_string("E5EAF2")  # text on navy
AMBER_BG = RGBColor.from_string("FEF6EC")
GREEN_BG = RGBColor.from_string("ECFDF3")
RED_BG   = RGBColor.from_string("FEF2F2")
TEAL_BG  = RGBColor.from_string("E6FAF7")

FONT   = "Be Vietnam Pro"          # body (bold=True → real Bold face)
FONTB  = "Be Vietnam Pro"          # titles/emphasis via bold flag (avoids faux-bold)
MONO   = "Consolas"                # code/diagram (Be Vietnam Pro has no mono)

# ---- geometry ---------------------------------------------------------------
EMU_IN = 914400
PW, PH = 13.333, 7.5
LM = 0.55
CW = PW - 2 * LM
CHECK, CROSS, STAR, ARROW, DOT = "✓", "✗", "★", "→", "●"

prs = Presentation()
prs.slide_width = Emu(int(PW * EMU_IN))
prs.slide_height = Emu(int(PH * EMU_IN))
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _noshadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def box(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines = list of dicts: text,size,color,bold,font,align,sa(space_after),sb,it"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if "sa" in ln:
            p.space_after = Pt(ln["sa"])
        if "sb" in ln:
            p.space_before = Pt(ln["sb"])
        runs = ln["runs"] if "runs" in ln else [ln]
        for rr in runs:
            r = p.add_run()
            r.text = rr["text"]
            f = r.font
            f.size = Pt(rr.get("size", ln.get("size", 16)))
            f.name = rr.get("font", ln.get("font", FONT))
            f.bold = rr.get("bold", ln.get("bold", False))
            f.italic = rr.get("it", ln.get("it", False))
            f.color.rgb = rr.get("color", ln.get("color", TEXT))
    return tb


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    _noshadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.text_frame.word_wrap = True
    sp.text_frame.margin_left = Inches(0.12)
    sp.text_frame.margin_right = Inches(0.12)
    sp.text_frame.margin_top = Inches(0.08)
    sp.text_frame.margin_bottom = Inches(0.08)
    return sp


def fill_shape_text(sp, lines, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if "sa" in ln:
            p.space_after = Pt(ln["sa"])
        if "sb" in ln:
            p.space_before = Pt(ln["sb"])
        runs = ln["runs"] if "runs" in ln else [ln]
        for rr in runs:
            r = p.add_run()
            r.text = rr["text"]
            f = r.font
            f.size = Pt(rr.get("size", ln.get("size", 14)))
            f.name = rr.get("font", ln.get("font", FONT))
            f.bold = rr.get("bold", ln.get("bold", False))
            f.italic = rr.get("it", ln.get("it", False))
            f.color.rgb = rr.get("color", ln.get("color", TEXT))


def title_bar(s, title, kicker=None, dark=False):
    tc = NAVYTX if dark else TEXT
    kc = TEAL
    if kicker:
        box(s, LM, 0.30, CW, 0.32, [{"text": kicker.upper(), "size": 12.5, "color": kc, "bold": True, "font": FONTB}])
        box(s, LM, 0.60, CW, 0.7, [{"text": title, "size": 27, "color": tc, "bold": True, "font": FONTB}])
        uy = 1.30
    else:
        box(s, LM, 0.42, CW, 0.7, [{"text": title, "size": 28, "color": tc, "bold": True, "font": FONTB}])
        uy = 1.18
    rect(s, LM, uy, 2.1, 0.055, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    return uy + 0.25


def footer(s, n, dark=False):
    c = RGBColor.from_string("7C8AA0") if dark else SUBTLE
    if not dark:
        rect(s, LM, 7.02, CW, 0.012, fill=BORDER, shape=MSO_SHAPE.RECTANGLE)
    box(s, LM, 7.06, CW * 0.6, 0.32,
        [{"text": "Neko Core  ·  Đại học Hàng hải Việt Nam (VMU)", "size": 9.5, "color": c}])
    box(s, LM + CW * 0.45, 7.06, CW * 0.55, 0.32,
        [{"text": f"neko-core:qwen3-4b-r2-cuda128-20260627   ·   {n}/14", "size": 9.5, "color": c, "align": PP_ALIGN.RIGHT}])


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def stat(s, l, t, w, number, label, color=TEAL, nsize=40, lsize=12):
    box(s, l, t, w, 0.85,
        [{"text": number, "size": nsize, "color": color, "bold": True, "font": FONTB, "align": PP_ALIGN.LEFT, "sa": 2}])
    box(s, l, t + 0.78, w, 0.5,
        [{"text": label, "size": lsize, "color": SUBTLE, "align": PP_ALIGN.LEFT}])


# ===========================================================================
# SLIDE 1 — Bìa
# ===========================================================================
s = slide(INK_BG)
# accent bar
rect(s, 0, 0, 0.16, PH, fill=TEAL, shape=MSO_SHAPE.RECTANGLE)
s.shapes.add_picture(BANNER, Inches(0.9), Inches(0.85), width=Inches(5.4))
box(s, 0.9, 2.0, 9.5, 1.2,
    [{"text": "Neko Core", "size": 58, "color": WHITE, "bold": True, "font": FONTB}])
box(s, 0.92, 3.15, 10.5, 0.6,
    [{"text": "Thuyết minh phương pháp", "size": 24, "color": TEAL, "bold": True, "font": FONTB}])
box(s, 0.92, 3.70, 10.5, 0.5,
    [{"text": "HackAIthon 2026 — Bảng C (Innovator)", "size": 16, "color": NAVYTX}])
# headline stat card (right)
c = rect(s, 9.35, 1.95, 3.4, 1.7, fill=RGBColor.from_string("121B2E"), line=TEAL, lw=1.25)
fill_shape_text(c, [
    {"text": "83.59", "size": 46, "color": TEAL, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER, "sa": 2},
    {"text": "public-463 leaderboard", "size": 12, "color": NAVYTX, "align": PP_ALIGN.CENTER, "sa": 4},
    {"text": "Qwen3-4B · offline · một mô hình ≤5B", "size": 10.5, "color": RGBColor.from_string("9FB0C7"), "align": PP_ALIGN.CENTER},
])
box(s, 0.92, 5.55, 11.6, 0.9, [
    {"text": "Đội Neko Core — Đại học Hàng hải Việt Nam (VMU)", "size": 14, "color": WHITE, "bold": True, "sa": 4},
    {"text": "Nguyễn Mạnh Hùng (CNT63ĐH, trưởng nhóm) · Bùi Việt Hoàng (CLC63ĐH) · Phạm Thị Minh Hồng (CNT63ĐH) · Phạm Thị Thu Thảo (KTN63ĐH) · Nghiêm Thị Mỹ Linh (KPM63ĐH)",
     "size": 11, "color": RGBColor.from_string("9FB0C7")},
])
notes(s, "Chào ban giám khảo. Chúng em là đội Neko Core, Đại học Hàng hải Việt Nam. Bài dự thi Bảng C của "
          "chúng em là một Docker image offline, tự chứa, chạy mô hình Qwen3-4B — đạt 83.59 trên leaderboard "
          "public. Nhưng điều chúng em muốn trình bày hôm nay không phải con số, mà là CÁCH chúng em ra quyết "
          "định để đạt nó.")

# ===========================================================================
# SLIDE 2 — Bài toán & hợp đồng chấm điểm
# ===========================================================================
s = slide()
y = title_bar(s, "Bài toán & hợp đồng chấm điểm", kicker="01 · Đề bài")
# left: contract flow
box(s, LM, y + 0.05, 6.0, 0.4, [{"text": "Hợp đồng I/O — Docker offline, tự chứa", "size": 14, "color": TEXT, "bold": True, "font": FONTB}])
flow = [("/code/private_test.json", SURF, TEXT), ("container offline", INK_BG, WHITE), ("/code/submission.csv", TEAL_BG, TEXT)]
fx = LM
for i, (lab, fl, tc) in enumerate(flow):
    w = 1.95
    sp = rect(s, fx, y + 0.55, w, 0.7, fill=fl, line=BORDER, lw=1.0)
    fill_shape_text(sp, [{"text": lab, "size": 11.5, "color": tc, "bold": True, "font": MONO, "align": PP_ALIGN.CENTER}])
    if i < 2:
        box(s, fx + w - 0.02, y + 0.55, 0.34, 0.7, [{"text": ARROW, "size": 20, "color": TEAL, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    fx += w + 0.32
box(s, LM, y + 1.45, 6.1, 1.0, [
    {"text": "Chấm trên 2000 câu private đa lĩnh vực", "size": 13, "color": TEXT, "bold": True, "sa": 3},
    {"text": "Kiểu THPT Quốc gia: Toán / Lý / Hoá / Sinh + Giáo dục công dân + đọc hiểu.", "size": 12, "color": SUBTLE, "sa": 6},
])
# ràng buộc card
cc = rect(s, LM, y + 2.55, 6.1, 1.55, fill=AMBER_BG, line=RGBColor.from_string("F3D9B0"), lw=1.0)
fill_shape_text(cc, [
    {"text": "Ràng buộc BTC (chính thức 2026-06-18)", "size": 12.5, "color": AMBER, "bold": True, "font": FONTB, "sa": 4},
    {"text": "• GPU chấm: RTX 5060 Ti (Blackwell, 16GB) · CUDA ≥ 12.8", "size": 12, "color": TEXT, "sa": 2},
    {"text": "• ≤ 5B TỔNG tham số · đúng 1 LLM · KHÔNG embedding/rerank", "size": 12, "color": TEXT, "sa": 2},
    {"text": "• Không API/mô hình ngoài · hoàn toàn offline", "size": 12, "color": TEXT},
], anchor=MSO_ANCHOR.TOP)
# right: scoring table
box(s, 7.05, y + 0.05, 5.7, 0.4, [{"text": "Trọng số chấm điểm", "size": 14, "color": TEXT, "bold": True, "font": FONTB}])
rows = [("Tiêu chí", "Điểm", "Hệ quả thiết kế", True),
        ("Accuracy", "80", "mục tiêu chính — lever phải tổng quát hoá", False),
        ("Thời gian inference", "10", "tốc độ là mục tiêu hạng nhất", False),
        ("Ý tưởng / sáng tạo", "10", "← tài liệu này", False)]
tbl = s.shapes.add_table(4, 3, Inches(7.05), Inches(y + 0.5), Inches(5.7), Inches(2.7)).table
tbl.columns[0].width = Inches(2.0); tbl.columns[1].width = Inches(0.9); tbl.columns[2].width = Inches(2.8)
for ri, (a, b, cc2, hdr) in enumerate(rows):
    for ci, val in enumerate((a, b, cc2)):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if hdr:
            cell.fill.fore_color.rgb = INK_BG
        elif a == "Ý tưởng / sáng tạo":
            cell.fill.fore_color.rgb = TEAL_BG
        else:
            cell.fill.fore_color.rgb = WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci == 1 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.size = Pt(13 if ci == 1 else 11.5)
        r.font.name = FONTB if (hdr or ci == 1) else FONT
        r.font.bold = hdr or ci == 1
        r.font.color.rgb = WHITE if hdr else (TEAL if (ci == 1 and not hdr) else TEXT)
footer(s, 2)
notes(s, "Đề bài: một container offline đọc đề ở /data, ghi đáp án ra /output. Chấm trên 2000 câu private. "
          "Accuracy chiếm 80 điểm, nhưng để ý hai cột còn lại — Time và Idea, mỗi cái 10 điểm — định hình "
          "toàn bộ thiết kế. Và có một ràng buộc cứng: tối đa 5 tỉ tham số, một mô hình, không công cụ ngoài.")

# ===========================================================================
# SLIDE 3 — GIẢI PHÁP (solution overview)
# ===========================================================================
s = slide()
y = title_bar(s, "Giải pháp: một Docker offline tự-chứa", kicker="02 · Giải pháp")
box(s, LM, y + 0.0, CW, 0.85, [{"runs": [
    {"text": "Neko Core đóng gói sẵn ", "size": 14.5, "color": TEXT},
    {"text": "Qwen3-4B (≤5B)", "size": 14.5, "color": TEAL, "bold": True, "font": FONTB},
    {"text": " trong một container offline: đọc /code/private_test.json → cho mô hình ", "size": 14.5, "color": TEXT},
    {"text": "LẬP LUẬN (CoT)", "size": 14.5, "color": TEXT, "bold": True, "font": FONTB},
    {"text": " → ghi /code/submission.csv — bọc trong các lớp ", "size": 14.5, "color": TEXT},
    {"text": "kỹ thuật chống-0-điểm", "size": 14.5, "color": GREEN, "bold": True, "font": FONTB},
    {"text": ".", "size": 14.5, "color": TEXT},
]}])
ry = y + 1.0
lft = rect(s, LM, ry + 0.25, 2.0, 1.4, fill=SURF, line=BORDER, lw=1.0)
fill_shape_text(lft, [{"text": "Đề thi", "size": 12.5, "color": TEXT, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER, "sa": 2},
                      {"text": "private_test.json", "size": 10, "color": SUBTLE, "font": MONO, "align": PP_ALIGN.CENTER},
                      {"text": "/code", "size": 10, "color": SUBTLE, "align": PP_ALIGN.CENTER}])
box(s, LM + 2.0, ry + 0.25, 0.32, 1.4, [{"text": ARROW, "size": 20, "color": TEAL, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
cen = rect(s, LM + 2.35, ry, 7.0, 1.9, fill=INK_BG, line=TEAL, lw=1.5)
fill_shape_text(cen, [
    {"text": "CONTAINER  (offline · ≤5B · 1 model)", "size": 12, "color": RGBColor.from_string("9FB0C7"), "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER, "sa": 5},
    {"text": "Qwen3-4B  +  chain-of-thought  (lập luận → trích đáp án)", "size": 14, "color": WHITE, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER, "sa": 7},
    {"text": "lớp chống-0-điểm:  arch-portable mọi GPU · contract submission.csv · checkpoint/resume", "size": 11, "color": TEAL, "align": PP_ALIGN.CENTER}])
box(s, LM + 9.35, ry + 0.25, 0.32, 1.4, [{"text": ARROW, "size": 20, "color": TEAL, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
rgt = rect(s, LM + 9.7, ry + 0.25, 2.53, 1.4, fill=TEAL_BG, line=TEAL, lw=1.0)
fill_shape_text(rgt, [{"text": "submission.csv", "size": 12, "color": TEXT, "bold": True, "font": MONO, "align": PP_ALIGN.CENTER, "sa": 2},
                      {"text": "+ submission_time.csv", "size": 9.5, "color": SUBTLE, "font": MONO, "align": PP_ALIGN.CENTER},
                      {"text": "trong /code", "size": 10, "color": SUBTLE, "align": PP_ALIGN.CENTER}])
gtees = [("≤5B · 1 model", "đúng luật BTC"), ("83.59", "public leaderboard"), ("~3 giây/câu", "đo thật trên GPU"), ("never-0", "chạy đúng mọi GPU")]
gw = (CW - 3 * 0.25) / 4
gx = LM
for big, sub in gtees:
    sp = rect(s, gx, ry + 2.25, gw, 1.0, fill=WHITE, line=BORDER, lw=1.0)
    fill_shape_text(sp, [{"text": big, "size": 17, "color": TEAL, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER, "sa": 2},
                         {"text": sub, "size": 10.5, "color": SUBTLE, "align": PP_ALIGN.CENTER}])
    gx += gw + 0.25
footer(s, 3)
notes(s, "Đây là giải pháp ở mức tổng quan. Neko Core là một Docker image offline, tự chứa, đã nướng sẵn mô "
          "hình Qwen3-4B bên trong. Nó đọc đề trắc nghiệm ở /data, cho mô hình lập luận bằng chain-of-thought "
          "rồi trích đáp án, và ghi ra /output/pred.csv. Bao quanh lõi đó là các lớp kỹ thuật chống-0-điểm: chạy "
          "đúng trên mọi GPU, hợp đồng pred.csv không thể hỏng, và checkpoint để không làm lại từ đầu. Kết quả: "
          "đúng luật ≤5B, đạt 83.59, khoảng 3 giây mỗi câu, và không bao giờ về 0 điểm.")

# ===========================================================================
# SLIDE 4 — Ý TƯỞNG CỐT LÕI (the one core idea)
# ===========================================================================
s = slide()
y = title_bar(s, "Ý tưởng cốt lõi", kicker="03 · Ý tưởng")
box(s, LM, y + 0.0, CW, 1.15, [
    {"runs": [
        {"text": "Đa số đội tối ưu accuracy ", "size": 17, "color": TEXT, "bold": True, "font": FONTB},
        {"text": "trên máy dev", "size": 17, "color": SUBTLE, "bold": True, "font": FONTB},
        {"text": ". Neko Core tối ưu ", "size": 17, "color": TEXT, "bold": True, "font": FONTB},
        {"text": "ĐIỂM KỲ VỌNG dưới bất định phần cứng", "size": 17, "color": TEAL, "bold": True, "font": FONTB},
        {"text": " — và để ", "size": 17, "color": TEXT, "bold": True, "font": FONTB},
        {"text": "PHÉP ĐO chọn giải pháp", "size": 17, "color": TEAL, "bold": True, "font": FONTB},
        {"text": ", không phải linh cảm.", "size": 17, "color": TEXT, "bold": True, "font": FONTB},
    ], "align": PP_ALIGN.CENTER}
], anchor=MSO_ANCHOR.MIDDLE)
cw2 = (CW - 0.4) / 2
ca = rect(s, LM, y + 1.35, cw2, 2.25, fill=WHITE, line=TEAL, lw=1.25)
fill_shape_text(ca, [
    {"text": "① Chọn-bằng-ĐO, không bằng linh cảm", "size": 14, "color": TEAL, "bold": True, "font": FONTB, "sa": 6},
    {"text": "Xây cả một 'xưởng' chiến lược → đo từng cái trên held-out → ship cái THẮNG (hoá ra là CoT đơn).", "size": 12, "color": TEXT, "sa": 6},
    {"text": "Loại thẳng tay cái hại — có bằng chứng đo:", "size": 12, "color": TEXT, "sa": 3},
    {"text": "FT −4.44 · TIR −16.52 · RAG ảo / phạm luật", "size": 12.5, "color": RED, "bold": True, "font": FONTB},
], anchor=MSO_ANCHOR.TOP)
cb = rect(s, LM + cw2 + 0.4, y + 1.35, cw2, 2.25, fill=WHITE, line=GREEN, lw=1.25)
fill_shape_text(cb, [
    {"text": "② Kỹ thuật KHÔNG-BAO-GIỜ-0-điểm", "size": 14, "color": GREEN, "bold": True, "font": FONTB, "sa": 6},
    {"text": "Thảm hoạ contest chấm-Docker = container 0 điểm (OOM / timeout / crash / sai GPU).", "size": 12, "color": TEXT, "sa": 6},
    {"text": "→ engineer: arch-portable mọi GPU + hợp đồng submission.csv bất khả xâm phạm + checkpoint/resume.", "size": 12, "color": TEXT, "sa": 4},
    {"text": "Robust trên MỌI phần cứng của BTC.", "size": 12, "color": TEXT, "bold": True},
], anchor=MSO_ANCHOR.TOP)
pc = rect(s, LM, y + 3.8, CW, 0.72, fill=INK_BG, line=None)
fill_shape_text(pc, [{"runs": [
    {"text": "Cái MỚI không nằm ở mô hình — mà ở ", "size": 15, "color": WHITE, "bold": True, "font": FONTB},
    {"text": "CÁCH chọn (đo lường)", "size": 15, "color": TEAL, "bold": True, "font": FONTB},
    {"text": " và ", "size": 15, "color": WHITE, "bold": True, "font": FONTB},
    {"text": "CÁCH đảm bảo (không-0-điểm)", "size": 15, "color": TEAL, "bold": True, "font": FONTB},
    {"text": ".", "size": 15, "color": WHITE, "bold": True, "font": FONTB},
], "align": PP_ALIGN.CENTER}])
footer(s, 4)
notes(s, "Đây là ý tưởng cốt lõi, nói thẳng. Đa số đội tối ưu accuracy trên máy dev của họ. Neko Core tối ưu "
          "điểm kỳ vọng dưới bất định phần cứng, và để phép đo chọn giải pháp chứ không phải linh cảm. Hai hệ "
          "quả: một, chúng em xây cả một xưởng chiến lược rồi đo từng cái trên held-out, ship cái thắng và loại "
          "thẳng cái hại — fine-tune, TIR, RAG đều bị loại bằng số đo. Hai, chúng em engineer để container không "
          "bao giờ về 0 điểm trên mọi GPU. Cái mới không ở mô hình mà ở cách chọn và cách đảm bảo.")

# ===========================================================================
# SLIDE 5 — Kiến trúc: foundry chiến lược
# ===========================================================================
s = slide()
y = title_bar(s, "Kiến trúc: một 'foundry chiến lược' composable", kicker="04 · Harness")
stages = ["loader", "router /\nclassifier", "solver\n(CoT · k=1)", "normalizer", "constrained\nrepair (GBNF)", "contract\nrepair"]
n = len(stages); sw = 1.78; sg = 0.16
total = n * sw + (n - 1) * sg
sx0 = (PW - total) / 2
sx = sx0
for i, st in enumerate(stages):
    fill = TEAL_BG if "solver" in st else WHITE
    line = TEAL if "solver" in st else BORDER
    sp = rect(s, sx, y + 0.2, sw, 0.85, fill=fill, line=line, lw=1.25 if "solver" in st else 1.0)
    fill_shape_text(sp, [{"text": st, "size": 11, "color": TEXT, "bold": "solver" in st, "font": MONO, "align": PP_ALIGN.CENTER}])
    if i < n - 1:
        box(s, sx + sw - 0.05, y + 0.2, 0.26, 0.85, [{"text": ARROW, "size": 16, "color": TEAL, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    sx += sw + sg
box(s, sx0, y + 1.15, total, 0.35, [{"text": "/code/private_test.json   " + ARROW + "   …   " + ARROW + "   /code/submission.csv + submission_time.csv", "size": 11.5, "color": SUBTLE, "font": MONO, "align": PP_ALIGN.CENTER}])
# swappable strategies row
box(s, LM, y + 1.7, CW, 0.35, [{"text": "Strategy hoán đổi được (typed processor, sau policy gate):", "size": 12.5, "color": TEXT, "bold": True}])
strat = [("self-consistency", True), ("tiered", False), ("TIR", False), ("reading", False), ("rag", False)]
sx = LM
for lab, ship in strat:
    w = 2.05
    sp = rect(s, sx, y + 2.1, w, 0.55, fill=TEAL_BG if ship else SURF, line=TEAL if ship else BORDER, lw=1.25 if ship else 1.0)
    fill_shape_text(sp, [{"text": (CHECK + " " if ship else "") + lab + ("  (ship)" if ship else ""), "size": 11.5, "color": GREEN if ship else SUBTLE, "bold": ship, "font": MONO, "align": PP_ALIGN.CENTER}])
    sx += w + 0.2
qc = rect(s, LM, y + 2.95, CW, 0.95, fill=INK_BG, line=None)
fill_shape_text(qc, [
    {"runs": [
        {"text": "Chúng em xây cả một FOUNDRY chiến lược — rồi chọn cấu hình ship ", "size": 15, "color": WHITE, "bold": True, "font": FONTB},
        {"text": "BẰNG ĐO held-out", "size": 15, "color": TEAL, "bold": True, "font": FONTB},
        {"text": ", không bằng linh cảm. Cái gì ship là kẻ thắng của một phép đo.", "size": 15, "color": WHITE, "bold": True, "font": FONTB},
    ], "align": PP_ALIGN.CENTER}
])
footer(s, 5)
notes(s, "Đây là kiến trúc. Provider giấu sau một hợp đồng complete() duy nhất; các strategy là processor "
          "gõ-kiểu, hoán đổi được qua một policy gate. Điểm mấu chốt: chúng em không đoán nên dùng chiến lược "
          "nào — chúng em xây tất cả, đo từng cái, rồi mới chọn.")

# ===========================================================================
# SLIDE 6 — Pivot luật bằng config
# ===========================================================================
s = slide()
y = title_bar(s, "Bước chuyển luật ≤5B: thích nghi bằng CONFIG", kicker="05 · Bản lĩnh thích nghi")
# before (left)
bc = rect(s, LM, y + 0.2, 5.6, 3.2, fill=SURF, line=BORDER, lw=1.0)
fill_shape_text(bc, [
    {"text": "TRƯỚC", "size": 12, "color": SUBTLE, "bold": True, "font": FONTB, "sa": 6},
    {"text": "Gemma-4-26B-A4B", "size": 19, "color": SUBTLE, "bold": True, "font": FONTB, "sa": 2, "it": True},
    {"text": "88.55 public-463 — nhưng 26B > 5B", "size": 13, "color": SUBTLE, "sa": 8},
    {"text": CROSS + "  BỊ LOẠI dưới luật mới", "size": 16, "color": RED, "bold": True, "font": FONTB},
], anchor=MSO_ANCHOR.MIDDLE)
box(s, 6.25, y + 1.35, 0.7, 0.8, [{"text": ARROW, "size": 34, "color": TEAL, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
# after (right)
ac = rect(s, 7.15, y + 0.2, 5.65, 3.2, fill=TEAL_BG, line=TEAL, lw=1.25)
fill_shape_text(ac, [
    {"text": "SAU — pivot = sửa MỘT khối config", "size": 12, "color": TEAL, "bold": True, "font": FONTB, "sa": 6},
    {"text": 'runtime.model_policy', "size": 14, "color": TEXT, "bold": True, "font": MONO, "sa": 2},
    {"text": '{ "aliases": ["*"], "max_params_b": 5.0 }', "size": 13, "color": INDIGO, "font": MONO, "sa": 8},
    {"text": "→ ép mọi mô hình ≤5B, KHÔNG sửa một dòng", "size": 13.5, "color": TEXT, "bold": True, "sa": 2},
    {"text": "code solver nào. Pivot sang Qwen3-4B.", "size": 13.5, "color": TEXT, "bold": True},
], anchor=MSO_ANCHOR.MIDDLE)
box(s, LM, y + 3.55, CW, 0.6, [{"runs": [
    {"text": "Đây là một đòn bẩy 'Ý tưởng': ", "size": 14, "color": TEXT},
    {"text": "harness thích nghi luật chơi bằng dữ liệu, không bằng code.", "size": 14, "color": TEAL, "bold": True, "font": FONTB},
], "align": PP_ALIGN.CENTER}])
footer(s, 6)
notes(s, "Giữa cuộc thi, ban tổ chức đổi luật từ mở-rộng sang tối đa 5 tỉ tham số. Mô hình cũ của chúng em, "
          "Gemma 26B, lập tức bị loại. Phần lớn đội sẽ phải sửa code. Chúng em chỉ sửa DỮ LIỆU config — "
          "allowlist mô hình là config-driven — và pivot sang Qwen3-4B mà không đụng một dòng code giải. Đó "
          "chính là 'ý tưởng & tối ưu sáng tạo': harness thích nghi luật bằng config.")

# ===========================================================================
# SLIDE 7 — Model ≤5B
# ===========================================================================
s = slide()
y = title_bar(s, "Mô hình thi đấu: Qwen3-4B-Instruct-2507", kicker="06 · Mô hình ≤5B")
mc = rect(s, LM, y + 0.2, 6.0, 3.4, fill=WHITE, line=BORDER, lw=1.0)
fill_shape_text(mc, [
    {"text": "Qwen3-4B-Instruct-2507", "size": 19, "color": TEXT, "bold": True, "font": FONTB, "sa": 8},
    {"text": CHECK + "  Dense 4B — ≤5B theo tổng tham số (hợp luật rõ ràng)", "size": 13, "color": TEXT, "sa": 5},
    {"text": CHECK + "  GGUF Q5_K_M · llama.cpp · hoàn toàn offline", "size": 13, "color": TEXT, "sa": 5},
    {"text": CHECK + "  ĐÚNG 1 mô hình — không embedding / rerank", "size": 13, "color": TEXT, "sa": 5},
    {"text": CHECK + "  Engine: chain-of-thought (CoT), k=1", "size": 13, "color": TEXT, "bold": True, "sa": 3},
    {"text": "      self-consistency voting k=5 đã đo — không cải thiện → giữ k=1", "size": 10.5, "color": SUBTLE},
], anchor=MSO_ANCHOR.MIDDLE)
# VRAM stat + bar
box(s, 7.05, y + 0.25, 5.7, 0.5, [{"text": "VRAM lúc chạy (đo thật)", "size": 13, "color": TEXT, "bold": True, "font": FONTB}])
stat(s, 7.05, y + 0.7, 5.7, "~5 GB", "trên RTX 5060 Ti 16GB → chừa ~11GB", color=TEAL, nsize=42)
# bar
rect(s, 7.05, y + 2.05, 5.5, 0.45, fill=SURF, line=BORDER, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
rect(s, 7.05, y + 2.05, 5.5 * 5 / 16, 0.45, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
box(s, 7.05, y + 2.55, 5.5, 0.35, [{"text": "5GB dùng  |  11GB dư  → không OOM trên mọi GPU thi đấu", "size": 11, "color": SUBTLE}])
footer(s, 7)
notes(s, "Mô hình thi đấu là Qwen3-4B-Instruct-2507, dense 4 tỉ tham số — dưới 5B rõ ràng. Chạy local qua "
          "llama.cpp, hoàn toàn offline, đúng một mô hình, không có công cụ phụ. Lúc chạy chỉ dùng khoảng 5GB "
          "trên 16GB — thừa sức. Engine cốt lõi là chain-of-thought ở k=1: cho mô hình lập luận rồi trích đáp án. "
          "Biến thể self-consistency voting k=5 đã đo nhưng không cải thiện nên giữ k=1.")

# ===========================================================================
# SLIDE 8 — Không-bao-giờ-0-điểm #1: arch-portability
# ===========================================================================
s = slide()
y = title_bar(s, "Không-bao-giờ-0-điểm #1: arch-portability", kicker="07 · Robustness · lớp 1")
box(s, LM, y + 0.1, CW, 0.4, [{"text": "GPU chấm = RTX 5060 Ti (Blackwell sm_120) · build CUDA 12.8 · PTX-portable mọi GPU ≥ sm_75:", "size": 13, "color": TEXT, "bold": True, "font": FONTB}])
gpus = [("sm_75", "T4"), ("sm_80", "A100"), ("sm_86", "Ampere"), ("sm_89", "Ada"), ("sm_90", "H100"), ("sm_120", "5060 Ti")]
gw = (CW - 5 * 0.12) / 6
gx = LM
for sm, name in gpus:
    sp = rect(s, gx, y + 0.6, gw, 0.9, fill=WHITE, line=TEAL, lw=1.0)
    fill_shape_text(sp, [
        {"text": sm, "size": 12, "color": TEAL, "bold": True, "font": MONO, "align": PP_ALIGN.CENTER, "sa": 1},
        {"text": name, "size": 10, "color": SUBTLE, "align": PP_ALIGN.CENTER}])
    gx += gw + 0.12
box(s, LM, y + 1.65, CW, 0.5, [{"runs": [
    {"text": "PTX ", "size": 13, "color": TEXT},
    {"text": "compute_75 + compute_120", "size": 13, "color": INDIGO, "bold": True, "font": MONO},
    {"text": "  → 5060 Ti JIT sm_120 lúc load (ngoài vòng đo Time).   ", "size": 13, "color": TEXT},
    {"text": "GGML_NATIVE=off", "size": 13, "color": INDIGO, "bold": True, "font": MONO},
    {"text": " → mọi CPU.", "size": 13, "color": TEXT},
]}])
vc = rect(s, LM, y + 2.35, CW, 1.05, fill=GREEN_BG, line=GREEN, lw=1.25)
fill_shape_text(vc, [
    {"text": CHECK + "  KIỂM CHỨNG", "size": 13.5, "color": GREEN, "bold": True, "font": FONTB, "sa": 4},
    {"text": "Smoke kiểu giám khảo (mount thư mục → /code) → predict.py đọc /code/private_test.json, ghi "
             "submission.csv + submission_time.csv hợp lệ. Blackwell sm_120 không máy nào test local được; "
             "image build CUDA 12.8 đúng yêu cầu BTC (BTC sẽ liên hệ nếu eval lỗi).",
     "size": 11.5, "color": TEXT}])
footer(s, 8)
notes(s, "Robustness lớp một: chúng em không biết ban tổ chức dùng GPU nào, nên nướng native SASS cho mọi kiến "
          "trúc NVIDIA từ 2016 tới 2025, kèm một PTX floor để mọi card mới hơn tự JIT lúc load. Build llama.cpp "
          "với GGML_NATIVE tắt để chạy được cả CPU đời cũ. Và chúng em đã verify thật: kéo nguyên image, kiểm "
          "tra đủ kiến trúc, chạy end-to-end ra pred.csv hợp lệ.")

# ===========================================================================
# SLIDE 9 — Không-bao-giờ-0-điểm #2: contract pred.csv
# ===========================================================================
s = slide()
y = title_bar(s, "Không-bao-giờ-0-điểm #2: hợp đồng submission.csv", kicker="08 · Robustness · lớp 2")
layers = [
    ("Mọi câu đều TRẢ LỜI ĐƯỢC", "Exception khi giải → fallback heuristic tất định; lỗi bất ngờ → fallback-after-error; chữ không hợp lệ → answer-repair."),
    ("submission.csv ghi TRƯỚC mọi thứ có thể raise", "Bước contract-repair phủ ĐÚNG mọi qid, mỗi câu 1 chữ cái hợp lệ theo số phương án TỪNG câu (giữ nguyên dự đoán tốt)."),
    ("Checkpoint mỗi câu + auto-resume", "Container bị ngắt giữa chừng vẫn tiếp tục từ checkpoint, không làm lại từ đầu."),
]
ly = y + 0.15
for i, (t, d) in enumerate(layers):
    sp = rect(s, LM, ly, CW, 0.92, fill=WHITE, line=BORDER, lw=1.0)
    nb = rect(s, LM + 0.2, ly + 0.21, 0.5, 0.5, fill=GREEN, shape=MSO_SHAPE.OVAL)
    fill_shape_text(nb, [{"text": CHECK, "size": 16, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}])
    box(s, LM + 0.95, ly + 0.13, CW - 1.2, 0.4, [{"text": t, "size": 14, "color": TEXT, "bold": True, "font": FONTB}])
    box(s, LM + 0.95, ly + 0.5, CW - 1.2, 0.4, [{"text": d, "size": 11.5, "color": SUBTLE}])
    ly += 1.04
qc = rect(s, LM, ly + 0.02, CW, 0.7, fill=INK_BG, line=None)
fill_shape_text(qc, [{"runs": [
    {"text": "Một lỗi bất kỳ ", "size": 15, "color": WHITE, "bold": True, "font": FONTB},
    {"text": "KHÔNG THỂ", "size": 15, "color": TEAL, "bold": True, "font": FONTB},
    {"text": " zero-hoá cả lần chạy 2000 câu.", "size": 15, "color": WHITE, "bold": True, "font": FONTB},
], "align": PP_ALIGN.CENTER}])
footer(s, 9)
notes(s, "Robustness lớp hai: hợp đồng file đầu ra. pred.csv được ghi TRƯỚC bất cứ bước nào có thể lỗi, và một "
          "bước contract-repair đảm bảo nó luôn phủ đúng mọi qid với một chữ cái hợp lệ. Cộng checkpoint mỗi "
          "câu và auto-resume nếu bị ngắt. Kết quả: một câu lỗi không bao giờ làm hỏng cả lần chạy 2000 câu.")

# ===========================================================================
# SLIDE 10 — Accuracy đo thật
# ===========================================================================
s = slide()
y = title_bar(s, "Hành trình Accuracy — đo thật", kicker="09 · Kết quả")
# left compare
box(s, LM, y + 0.1, 6.0, 0.4, [{"text": "Cấu hình đã đo", "size": 13.5, "color": TEXT, "bold": True, "font": FONTB}])
r1 = rect(s, LM, y + 0.55, 6.0, 0.7, fill=SURF, line=BORDER, lw=1.0)
fill_shape_text(r1, [{"runs": [{"text": "Letter-only (ép 1 token)  ", "size": 12.5, "color": SUBTLE}, {"text": "→  thấp (lập luận bị nén)", "size": 12.5, "color": SUBTLE, "bold": True}]}])
r2 = rect(s, LM, y + 1.35, 6.0, 0.85, fill=TEAL_BG, line=TEAL, lw=1.25)
fill_shape_text(r2, [
    {"runs": [{"text": CHECK + " Chain-of-thought k=1 (ship)", "size": 13, "color": TEXT, "bold": True, "font": FONTB}], "sa": 2},
    {"runs": [{"text": "80.22 proxy-450", "size": 13, "color": GREEN, "bold": True}, {"text": "   ·   ", "size": 13, "color": SUBTLE}, {"text": "83.59 leaderboard-463", "size": 13, "color": GREEN, "bold": True}]},
])
# cluster stats
box(s, 7.05, y + 0.1, 5.7, 0.4, [{"text": "Phân rã cụm (proxy-450)", "size": 13.5, "color": TEXT, "bold": True, "font": FONTB}])
cl = [("Toán/Hoá", "73.91"), ("GDCD", "78.67"), ("Khoa học", "85.41")]
cx = 7.05
for name, v in cl:
    sp = rect(s, cx, y + 0.55, 1.78, 1.0, fill=WHITE, line=BORDER, lw=1.0)
    fill_shape_text(sp, [
        {"text": v, "size": 22, "color": INDIGO, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER, "sa": 2},
        {"text": name, "size": 11, "color": SUBTLE, "align": PP_ALIGN.CENTER}])
    cx += 1.96
fc = rect(s, LM, y + 2.6, CW, 1.1, fill=AMBER_BG, line=RGBColor.from_string("F3D9B0"), lw=1.0)
fill_shape_text(fc, [
    {"text": "PHÁT HIỆN THEN CHỐT", "size": 12, "color": AMBER, "bold": True, "font": FONTB, "sa": 4},
    {"runs": [
        {"text": "Lỗi của một mô hình 4B trên đề THPT chủ yếu là ", "size": 14, "color": TEXT},
        {"text": "SAI SUY LUẬN / tính toán", "size": 14, "color": RED, "bold": True, "font": FONTB},
        {"text": ", không phải thiếu fact tra-cứu-được.", "size": 14, "color": TEXT},
    ], "sa": 2},
    {"text": "→ định hướng trực tiếp việc chọn / loại đòn bẩy ở slide sau.", "size": 12.5, "color": SUBTLE, "it": True},
])
footer(s, 10)
notes(s, "Về accuracy: ép mô hình trả một token thì điểm thấp vì lập luận bị nén. Cho nó lập luận đầy đủ rồi "
          "trích chữ cái (CoT) — nhảy lên 83.59 trên leaderboard. Phân rã theo cụm cho một phát "
          "hiện then chốt: lỗi của một mô hình 4B chủ yếu là sai suy luận, chứ không phải thiếu kiến thức "
          "tra-cứu-được. Phát hiện này quyết định chúng em chọn — và loại — lever nào.")

# ===========================================================================
# SLIDE 11 — ⭐ Killed levers (star)
# ===========================================================================
s = slide()
y = title_bar(s, STAR + " Kỷ luật chống overfit: BẰNG CHỨNG ĐO", kicker="10 · Bằng chứng đo lường")
box(s, LM, y + 0.0, CW, 0.45, [{"runs": [
    {"text": "Mỗi đòn bẩy hào nhoáng đều được ", "size": 13.5, "color": TEXT},
    {"text": "XÂY rồi LOẠI BỎ", "size": 13.5, "color": RED, "bold": True, "font": FONTB},
    {"text": " bằng một phép đo held-out — không phải linh cảm:", "size": 13.5, "color": TEXT},
]}])
data = [
    ("Đòn bẩy (đã xây + đo)", "Δ accuracy (Qwen3-4B, proxy)", "Quyết định", None),
    ("Fine-tune v1 (math+legal+mcq)", "−4.44  (quant −14.78)", CROSS + " LOẠI", RED),
    ("Fine-tune v2 (hết data in-dist)", "±0.00  (đúng base)", CROSS + " LOẠI", RED),
    ("RAG-gated", "+3.11 ẢO (variance; civics −5 sạch)", CROSS + " LOẠI (+ phạm luật)", RED),
    ("TIR (model viết + chạy Python)", "−16.52  (41% câu degrade)", CROSS + " LOẠI", RED),
    ("k5-vote / tiered", "≈ ship k=1 (trong nhiễu)", DOT + " không bật", SUBTLE),
]
tbl = s.shapes.add_table(len(data), 3, Inches(LM), Inches(y + 0.55), Inches(CW), Inches(3.0)).table
tbl.columns[0].width = Inches(4.6); tbl.columns[1].width = Inches(4.9); tbl.columns[2].width = Inches(2.73)
for ri, (a, b, c, col) in enumerate(data):
    hdr = ri == 0
    for ci, val in enumerate((a, b, c)):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK_BG if hdr else (RED_BG if (col == RED) else WHITE)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.size = Pt(11.5 if not hdr else 11.5)
        r.font.name = FONTB if (hdr or ci == 2) else FONT
        r.font.bold = hdr or ci == 2
        if hdr:
            r.font.color.rgb = WHITE
        elif ci == 1 and col == RED:
            r.font.color.rgb = RED
        elif ci == 2:
            r.font.color.rgb = col if col else TEXT
        else:
            r.font.color.rgb = TEXT
qc = rect(s, LM, y + 3.75, CW, 0.62, fill=GREEN_BG, line=GREEN, lw=1.25)
fill_shape_text(qc, [{"runs": [
    {"text": "Ship = CoT, k=1.  ", "size": 14, "color": GREEN, "bold": True, "font": FONTB},
    {"text": "KHÔNG thêm phức tạp gây hại — cũng là một dạng tối ưu.", "size": 14, "color": TEXT, "bold": True, "font": FONTB},
], "align": PP_ALIGN.CENTER}])
footer(s, 11)
notes(s, "Đây là phần chúng em tự hào nhất. Mỗi đòn bẩy hào nhoáng về mặt kỹ thuật — fine-tune, RAG, mô hình tự viết code Python — "
          "chúng em đều XÂY thật rồi LOẠI BỎ bằng một phép đo held-out. Fine-tune làm tệ đi 4.44 điểm vì "
          "catastrophic forgetting. RAG chỉ là ảo giác của variance, và nay còn phạm luật. TIR phá 16 điểm vì "
          "4B chưa đủ sức viết code đúng. Bài học: trên một base ≤5B đã mạnh, thêm phức tạp sai cách làm TỆ ĐI. "
          "Việc KHÔNG thêm thứ gây hại, bản thân nó, là một dạng tối ưu.")

# ===========================================================================
# SLIDE 12 — Time: minimalism = nhanh nhất
# ===========================================================================
s = slide()
y = title_bar(s, "Tối ưu Time: minimalism = nhanh nhất", kicker="11 · Điểm Time (Vòng-2)")
lc = rect(s, LM, y + 0.2, 6.1, 3.1, fill=WHITE, line=BORDER, lw=1.0)
fill_shape_text(lc, [
    {"text": "Đường ship tối giản = đường NHANH nhất", "size": 15, "color": TEXT, "bold": True, "font": FONTB, "sa": 8},
    {"text": "1 model + self-consistency.", "size": 13, "color": TEXT, "sa": 8},
    {"runs": [{"text": CROSS + " TIR ", "size": 12.5, "color": RED, "bold": True}, {"text": "(2 round + chạy code) · ", "size": 12.5, "color": SUBTLE}, {"text": CROSS + " RAG ", "size": 12.5, "color": RED, "bold": True}, {"text": "(retrieval) · ", "size": 12.5, "color": SUBTLE}, {"text": CROSS + " k>1 ", "size": 12.5, "color": RED, "bold": True}, {"text": "(nhiều mẫu)", "size": 12.5, "color": SUBTLE}], "sa": 6},
    {"text": "đều chậm hơn — mà đo ra đều KHÔNG giúp accuracy.", "size": 13, "color": TEXT, "bold": True},
], anchor=MSO_ANCHOR.MIDDLE)
# right stats
stat(s, 7.4, y + 0.35, 5.3, "~3 giây", "mỗi câu (đo thật trên GPU)", color=TEAL, nsize=40)
stat(s, 7.4, y + 1.55, 5.3, "~100 phút", "cho 2000 câu private", color=INDIGO, nsize=40)
gc = rect(s, 7.4, y + 2.7, 5.3, 0.6, fill=GREEN_BG, line=GREEN, lw=1.0)
fill_shape_text(gc, [{"text": CHECK + "  VRAM ~5GB → không OOM, không timeout (RTX 5060 Ti 16GB)", "size": 11.5, "color": TEXT, "bold": True, "align": PP_ALIGN.CENTER}])
box(s, LM, y + 3.5, CW, 0.5, [{"text": "Minimalism thắng CẢ HAI mặt trận: Accuracy (không hồi quy) và Time.", "size": 14, "color": TEAL, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER}])
footer(s, 12)
notes(s, "Điều đẹp nhất: đường chạy tối giản của chúng em cũng là đường nhanh nhất. Mọi lever phức tạp vừa "
          "không giúp accuracy vừa chậm hơn. Minimalism thắng cả hai mặt trận. Đo thật: ba giây một câu, 2000 "
          "câu khoảng 100 phút, 5GB VRAM — không bao giờ OOM hay timeout trên server 16GB.")

# ===========================================================================
# SLIDE 13 — Định vị trung thực
# ===========================================================================
s = slide()
y = title_bar(s, "Định vị trung thực", kicker="12 · Sự trưởng thành")
box(s, 1.0, y + 0.35, PW - 2.0, 1.0, [{"runs": [
    {"text": "Qwen3-4B đã gần ", "size": 23, "color": TEXT, "bold": True, "font": FONTB},
    {"text": "TRẦN ≤5B (~80–84)", "size": 23, "color": TEAL, "bold": True, "font": FONTB},
    {"text": " trên phân phối đề này.", "size": 23, "color": TEXT, "bold": True, "font": FONTB},
], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.4, y + 1.55, PW - 2.8, 1.3, [
    {"text": "Cho thấy bằng dữ liệu: mọi lever còn lại (FT / RAG / TIR / k5) đều nằm trong nhiễu hoặc gây hại "
             "trên đo sạch — đây là kết luận có dữ liệu, không phải bỏ cuộc.", "size": 14, "color": TEXT, "align": PP_ALIGN.CENTER, "sa": 8},
    {"runs": [
        {"text": "Đường tăng tiếp theo (nếu có) = ", "size": 14, "color": TEXT, "align": PP_ALIGN.CENTER},
        {"text": "thêm data MCQ đa môn", "size": 14, "color": INDIGO, "bold": True, "font": FONTB},
        {"text": " — bài toán thu thập dữ liệu, không phải kỹ thuật runtime.", "size": 14, "color": TEXT},
    ], "align": PP_ALIGN.CENTER},
])
bc = rect(s, 3.2, y + 3.0, PW - 6.4, 0.75, fill=INK_BG, line=None)
fill_shape_text(bc, [{"text": "Chúng em báo số đo thật — không hứa số.", "size": 16, "color": TEAL, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER}])
box(s, 1.4, y + 3.95, PW - 2.8, 0.85, [{"runs": [
    {"text": "Điểm yếu đã đo thẳng thắn: ", "size": 12.5, "color": SUBTLE},
    {"text": "cụm Sử/Địa (~62–65)", "size": 12.5, "color": AMBER, "bold": True, "font": FONTB},
    {"text": " — lỗi chủ yếu là chọn phương án 'chính' theo quy ước, không phải thiếu fact (nên RAG/FT không cứu được).", "size": 12.5, "color": SUBTLE},
], "align": PP_ALIGN.CENTER}])
footer(s, 13)
notes(s, "Chúng em định vị trung thực: Qwen3-4B đã gần trần của lớp dưới-5B trên phân phối đề này — và chúng em "
          "chứng minh điều đó bằng dữ liệu, chứ không phải bỏ cuộc. Đường tăng điểm tiếp theo là thêm dữ liệu, "
          "một bài toán thu thập, không phải kỹ thuật runtime. Mọi con số trong bài này là số đo thật.")

# ===========================================================================
# SLIDE 14 — Kết: 3 trụ Idea
# ===========================================================================
s = slide(INK_BG)
rect(s, 0, 0, 0.16, PH, fill=TEAL, shape=MSO_SHAPE.RECTANGLE)
box(s, 0.9, 0.55, CW, 0.4, [{"text": "TỔNG KẾT — BA TRỤ Ý TƯỞNG", "size": 14, "color": TEAL, "bold": True, "font": FONTB}])
pillars = [
    ("1", "Harness composable", "Một foundry chiến lược, gõ-kiểu, hoán đổi được — không phải mớ code dán."),
    ("2", "Chọn-lever-bằng-ĐO", "Bằng chứng đo held-out (FT −4.44 · TIR −16.52 · RAG ảo) — kỷ luật, không linh cảm."),
    ("3", "Kỹ thuật chống-0-điểm", "Arch-portable mọi GPU + hợp đồng submission.csv bất khả xâm phạm."),
]
pw_ = 3.9; gap = 0.35; total = 3 * pw_ + 2 * gap; px = (PW - total) / 2
for num, t, d in pillars:
    sp = rect(s, px, 1.5, pw_, 2.5, fill=RGBColor.from_string("121B2E"), line=TEAL, lw=1.0)
    nb = rect(s, px + pw_ / 2 - 0.4, 1.75, 0.8, 0.8, fill=TEAL, shape=MSO_SHAPE.OVAL)
    fill_shape_text(nb, [{"text": num, "size": 26, "color": INK_BG, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER}])
    box(s, px + 0.2, 2.75, pw_ - 0.4, 0.7, [{"text": t, "size": 16, "color": WHITE, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER}])
    box(s, px + 0.25, 3.35, pw_ - 0.5, 0.6, [{"text": d, "size": 11.5, "color": RGBColor.from_string("9FB0C7"), "align": PP_ALIGN.CENTER}])
    px += pw_ + gap
# đối chiếu thẳng rubric Ý tưởng (giúp giám khảo tick đúng ô chấm)
box(s, 1.0, 4.18, PW - 2.0, 0.34, [{"text": "ĐỐI CHIẾU RUBRIC Ý TƯỞNG (10đ) — phương pháp tối ưu: sáng tạo · hiệu quả · khả thi", "size": 12.5, "color": RGBColor.from_string("9FB0C7"), "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER}])
rubric = [
    ("✓ SÁNG TẠO", "harness foundry · chọn-lever-bằng-đo · pivot luật bằng config"),
    ("✓ HIỆU QUẢ", "83.59 leaderboard · ~3 giây/câu (nhanh nhất)"),
    ("✓ KHẢ THI", "offline ≤5GB · never-0 · CUDA 12.8 · Blackwell"),
]
px = (PW - total) / 2
for head, sub in rubric:
    sp = rect(s, px, 4.55, pw_, 0.92, fill=RGBColor.from_string("121B2E"), line=TEAL, lw=1.0)
    fill_shape_text(sp, [
        {"text": head, "size": 14, "color": TEAL, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER, "sa": 2},
        {"text": sub, "size": 9.5, "color": RGBColor.from_string("9FB0C7"), "align": PP_ALIGN.CENTER}])
    px += pw_ + gap
box(s, 1.0, 5.75, PW - 2.0, 0.55, [{"text": "Cảm ơn ban giám khảo đã lắng nghe.", "size": 18, "color": WHITE, "bold": True, "font": FONTB, "align": PP_ALIGN.CENTER}])
box(s, 1.0, 6.5, PW - 2.0, 0.4, [{"text": "image: hacamy12345/neko-core:qwen3-4b-r2-cuda128-20260627   ·   Neko Core — VMU", "size": 11, "color": RGBColor.from_string("7C8AA0"), "align": PP_ALIGN.CENTER}])

prs.save(OUT)
print("OK ->", OUT)
print("slides:", len(prs.slides._sldIdLst))
