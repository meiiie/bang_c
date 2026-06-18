"""Build the Vietnamese method-writeup slide deck for HackAIthon 2026 Bang C.

Reproducible: `python scripts/build_pptx.py` -> docs/Neko-Core-Thuyet-minh-phuong-phap.pptx

Layout follows the assertion-evidence method (Michael Alley): each slide headline is a
full sentence, the body is visual evidence. Copy register: academic / technical Vietnamese
(measured tone, precise terminology). Content mirrors docs/method-writeup-vi.md
(v0.7.2: Qwen3-4B-Instruct-2507, dense <=5B, leaderboard 83.59). python-pptx only.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- brand palette ----
ORANGE = RGBColor(0xF9, 0x73, 0x16)
ORANGE_DK = RGBColor(0xC2, 0x55, 0x0B)
CREAM = RGBColor(0xFA, 0xF5, 0xEE)
INK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x15, 0x80, 0x3D)
RED = RGBColor(0xB9, 0x1C, 0x1C)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
CODEBG = RGBColor(0x11, 0x18, 0x27)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = int(Inches(13.333))
prs.slide_height = int(Inches(7.5))
BLANK = prs.slide_layouts[6]


def _set(run, size, color=INK, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def box(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=3.0, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        if isinstance(para, tuple):
            para = [para]
        for seg in para:
            s, size, color, bold, italic = (seg + (INK, False, False))[:5]
            r = p.add_run(); r.text = s; _set(r, size, color, bold, italic)
    return tb


def header(slide, kicker, assertion):
    box(slide, 0, 0, 13.333, 7.5, fill=WHITE)
    box(slide, 0, 0, 13.333, 0.16, fill=ORANGE)
    text(slide, 0.7, 0.46, 12.0, 0.35, [[(kicker, 11, ORANGE, True, False)]])
    text(slide, 0.7, 0.78, 12.2, 1.1, [[(assertion, 22, INK, True, False)]])


def table(slide, x, y, w, rows, col_w, header_fill=INK, font_sz=12, head_sz=12, row_h=0.42):
    nrows, ncols = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(row_h * nrows)).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.1); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt, *style = (cell if isinstance(cell, tuple) else (cell,))
            color = style[0] if style else (WHITE if i == 0 else INK)
            bold = style[1] if len(style) > 1 else (i == 0)
            c.fill.solid()
            c.fill.fore_color.rgb = header_fill if i == 0 else (LIGHT if i % 2 == 0 else WHITE)
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(txt)
            _set(r, head_sz if i == 0 else font_sz, color, bold)
    return gt


def page_no(slide, n):
    text(slide, 12.4, 7.05, 0.8, 0.3, [[(str(n), 10, GRAY, False, False)]], align=PP_ALIGN.RIGHT)


# ============================ 1 · COVER ============================
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, fill=INK)
box(s, 0, 0, 13.333, 0.9, fill=ORANGE)
box(s, 0, 6.7, 13.333, 0.8, fill=ORANGE)
text(s, 0.9, 2.0, 11.5, 1.0, [[("NEKO CORE", 52, WHITE, True, False)]])
text(s, 0.9, 3.1, 11.6, 0.9, [[("Hệ thống suy luận trắc nghiệm ngoại tuyến ≤5B — tối ưu để không bao giờ ăn 0 điểm, mọi quyết định chứng minh bằng đo lường", 17, CREAM, False, False)]])
text(s, 0.9, 4.25, 11.5, 0.6, [[("Tài liệu thuyết minh phương pháp · HackAIthon 2026 — Bảng C (Innovator)", 16, ORANGE, True, False)]])
text(s, 0.9, 5.2, 11.5, 1.1, [
    [("Public‑463 leaderboard: ", 15, CREAM, False, False), ("83.59", 18, WHITE, True, False),
     ("   ·   Mô hình: Qwen3‑4B‑Instruct‑2507 (dense 4B, ≤5B, ngoại tuyến)", 14, CREAM, False, False)],
    [("Đội Neko Core · Đại học Hàng hải Việt Nam (VMU)", 12, RGBColor(0xC9, 0xCE, 0xD6), False, True)],
])

# ============================ 2 · TÓM TẮT ============================
s = prs.slides.add_slide(BLANK)
header(s, "Tóm tắt", "Một mô hình ≤5B, tối ưu trên ba trụ cột: suy luận đo được, không‑0‑điểm, chọn kỹ thuật bằng đo.")
cards = [
    ("Một mô hình ≤5B, suy luận tự nhất quán", "Qwen3‑4B‑Instruct‑2507 (dense 4B) chạy ngoại tuyến, chain‑of‑thought tự nhất quán, đạt 83.59 trên leaderboard. Một mô hình duy nhất — không mô hình phụ, không embedding/rerank.", ORANGE),
    ("Không bao giờ 0 điểm trên mọi máy", "Wheel biên dịch native cho mọi GPU NVIDIA 2016→2025 (P100→Blackwell) kèm PTX floor, chạy mọi CPU, hợp đồng pred.csv bốn lớp phục hồi. Dùng ~5GB/16GB VRAM. Đã kiểm chứng literal.", GREEN),
    ("Chọn kỹ thuật bằng đo, không đoán", "Đã xây cả một foundry chiến lược (fine‑tune, RAG, TIR, tiered) rồi loại thẳng cái gây hại bằng đo held‑out. Self‑consistency trên base mạnh là tối ưu cho ≤5B và đồng thời nhanh nhất.", INK),
]
for i, (t, d, c) in enumerate(cards):
    cx = 0.7 + i * 4.07
    ch = 3.15
    tc = c if c != INK else ORANGE_DK
    box(s, cx, 2.25, 3.85, ch, fill=WHITE, line=c, line_w=1.5)
    box(s, cx, 2.25, 3.85, 0.16, fill=c)
    text(s, cx + 0.32, 2.45, 3.25, ch - 0.35,
         [[(t, 15.5, tc, True, False)], [("", 8, INK, False, False)], [(d, 12, INK, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space=4.0)
text(s, 0.7, 6.05, 12, 0.7, [[("Định vị khách quan: 83.59 tiệm cận giới hạn thực tế của một mô hình ≤5B trên phân phối đề này; mọi đòn bẩy còn lại đều nằm trong sai số hoặc gây hại khi đo sạch. Tài liệu trình bày số liệu đo được, không cam kết số liệu chưa kiểm chứng.", 12, GRAY, False, True)]])
page_no(s, 2)

# ============================ 3 · BÀI TOÁN ============================
s = prs.slides.add_slide(BLANK)
header(s, "01 · Bài toán & hợp đồng", "Container ngoại tuyến đọc /data và ghi /output/pred.csv; chấm theo Độ chính xác 80, Thời gian 10, Ý tưởng 10.")
text(s, 0.7, 2.2, 7.3, 2.4, [
    [("Hệ thống đóng gói trong một container ", 13.5, INK, False, False),
     ("ngoại tuyến, tự chứa", 13.5, ORANGE_DK, True, False),
     (" (không mạng, không khóa API), đọc ", 13.5, INK, False, False),
     ("/data/*_test.csv", 13.5, ORANGE_DK, True, False), (" và ghi ", 13.5, INK, False, False),
     ("/output/pred.csv", 13.5, ORANGE_DK, True, False), (" với hai cột qid, answer.", 13.5, INK, False, False)],
    [("", 7, INK, False, False)],
    [("Tập kiểm thử riêng 2000 câu kiểu THPT Quốc gia (Toán/Lý/Hoá/Sinh, Giáo dục công dân, đọc hiểu); mọi kỹ thuật gắn cứng với một câu hay một ngôn ngữ đều bị loại hoặc tổng quát hóa qua cấu hình.", 13.5, INK, False, False)],
])
table(s, 8.3, 2.3, 4.4, [
    [("Tiêu chí Vòng 2", WHITE, True), ("Điểm", WHITE, True)],
    [("Độ chính xác", INK, True), ("80",)],
    [("Thời gian inference", INK, True), ("10",)],
    [("Ý tưởng / tối ưu", ORANGE_DK, True), ("10",)],
], [3.1, 1.3], row_h=0.55, font_sz=13, head_sz=13)
box(s, 0.7, 5.0, 12.0, 1.55, fill=CREAM, line=ORANGE, line_w=1.25)
text(s, 0.95, 5.15, 11.6, 1.3, [
    [("Ràng buộc Ban Tổ chức (xác nhận chính thức 2026‑06‑18):", 13, ORANGE_DK, True, False)],
    [("• Server chấm 16GB VRAM.   • ≤5B tính theo TỔNG tham số (nên MoE 26B‑tổng bị loại; dense Qwen3‑4B = 4B hợp lệ).", 12.5, INK, False, False)],
    [("• Chỉ 1 model LLM ≤5B, KHÔNG embedding/rerank, không mô hình/API ngoài, ngoại tuyến.  →  Bài nộp tuân thủ tuyệt đối.", 12.5, INK, False, False)],
])
page_no(s, 3)

# ============================ 4 · NGUYÊN TẮC ============================
s = prs.slides.add_slide(BLANK)
header(s, "02 · Nguyên tắc thiết kế", "Mô hình suy luận; hệ thống điều phối, đo lường và bảo vệ kết quả khỏi 0 điểm.")
items = [
    ("Phân tách vai trò", "Không mã hóa cứng đáp án hay công thức; mô hình thực hiện suy luận, hệ thống đo độ tin cậy từ mức đồng thuận của các mẫu."),
    ("Chống quá khớp bằng đo lường", "Tập riêng 2000 câu là mục tiêu, proxy 463 chỉ là que thăm. Một kỹ thuật chỉ triển khai khi tổng quát hóa được — đo held‑out, không hồi quy cụm."),
    ("Bền vững là ưu tiên hàng đầu", "Trong contest chấm‑Docker, rủi ro thật là container ăn 0 điểm (tràn bộ nhớ, hết giờ, sai GPU). Tối ưu điểm kỳ vọng dưới bất định phần cứng."),
    ("Ranh giới triển khai", "Container ngoại tuyến, tự chứa; các công cụ truy vết, đánh giá, phân tích chỉ tồn tại trong môi trường phát triển."),
]
for i, (t, d) in enumerate(items):
    cx = 0.7 + (i % 2) * 6.15
    cy = 2.25 + (i // 2) * 2.2
    box(s, cx, cy, 5.85, 1.95, fill=CREAM)
    box(s, cx, cy, 0.1, 1.95, fill=ORANGE)
    text(s, cx + 0.35, cy + 0.2, 5.3, 1.6,
         [[(f"{i+1}. {t}", 15, ORANGE_DK, True, False)], [("", 6, INK, False, False)], [(d, 12.5, INK, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space=4.0)
page_no(s, 4)

# ============================ 5 · KIẾN TRÚC ============================
s = prs.slides.add_slide(BLANK)
header(s, "03 · Kiến trúc", "Đường ống cấu hình hóa: bổ sung môn học, dạng câu hỏi hay đổi luật chỉ là thay đổi cấu hình.")
flow = ["/data CSV·JSON", "loader", "router / classifier", "solver — self‑consistency CoT",
        "answer normalizer", "constrained + contract repair", "/output/pred.csv"]
for i, step in enumerate(flow):
    cy = 2.2 + i * 0.62
    hot = i in (3, 6)
    box(s, 0.9, cy, 6.0, 0.5, fill=ORANGE if hot else CREAM)
    text(s, 1.15, cy + 0.06, 5.6, 0.4, [[(step, 13, WHITE if hot else INK, hot, False)]])
    if i < len(flow) - 1:
        text(s, 3.55, cy + 0.46, 1, 0.2, [[("▼", 9, ORANGE, True, False)]])
box(s, 7.5, 2.2, 5.2, 3.95, fill=LIGHT)
box(s, 7.5, 2.2, 5.2, 0.12, fill=INK)
text(s, 7.78, 2.45, 4.7, 3.6, [
    [("Foundry chiến lược, chọn bằng đo", 14, ORANGE_DK, True, False)],
    [("Các chiến lược (self‑consistency, tiered, TIR, reading) là thành phần gõ kiểu, hoán đổi được; cấu hình ship là kẻ thắng của một phép đo held‑out.", 12, INK, False, False)],
    [("Allowlist mô hình config‑driven: đổi luật Ban Tổ chức (26B → ≤5B) chỉ là sửa dữ liệu runtime.model_policy, không sửa logic.", 12, INK, False, False)],
    [("Cổng kiểm soát chính sách ngăn mọi thành phần chỉ phục vụ phát triển lọt vào ảnh nộp.", 12, INK, False, False)],
])
page_no(s, 5)

# ============================ 6 · ROBUSTNESS / ARCH (lá bài mạnh) ============================
s = prs.slides.add_slide(BLANK)
header(s, "04 · Bền vững phần cứng", "Một wheel chạy mọi GPU NVIDIA 2016→2025 — bài nộp không thể ăn 0 điểm vì sai phần cứng.")
text(s, 0.7, 2.15, 12.0, 0.6, [[
    ("Mô hình ", 13.5, INK, False, False), ("Qwen3‑4B Q5_K_M (~2.7GB), dùng ~5GB/16GB VRAM", 13.5, ORANGE_DK, True, False),
    (" — thừa sức trên server 16GB của Ban Tổ chức, không rủi ro tràn bộ nhớ.", 13.5, INK, False, False)]])
archs = [("sm_60", "P100"), ("sm_70", "V100"), ("sm_75", "T4"), ("sm_80", "A100"),
         ("sm_86", "Ampere"), ("sm_89", "Ada"), ("sm_90", "H100"), ("sm_120", "Blackwell")]
for i, (a, g) in enumerate(archs):
    cx = 0.7 + (i % 4) * 3.05
    cy = 2.95 + (i // 4) * 1.0
    box(s, cx, cy, 2.85, 0.82, fill=CREAM, line=ORANGE, line_w=1.0)
    text(s, cx, cy + 0.08, 2.85, 0.66, [[(a, 14, ORANGE_DK, True, False)], [(g, 11, GRAY, False, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=1.0)
text(s, 0.7, 5.15, 12.0, 1.4, [
    [("Native SASS cho 8 kiến trúc trên ", 12.5, INK, False, False), ("+ PTX floor compute_60", 12.5, GREEN, True, False),
     (" → mọi GPU ≥ Pascal tự JIT lúc nạp. ", 12.5, INK, False, False),
     ("GGML_NATIVE=off", 12.5, ORANGE_DK, True, False), (" → chạy mọi CPU (không SIGILL trên CPU cũ).", 12.5, INK, False, False)],
    [("Đã kiểm chứng LITERAL (2026‑06‑18): kéo nguyên ảnh từ Docker Hub trên GPU sạch → cuobjdump xác nhận đủ sm_60…sm_120 trong wheel → chroot chạy entrypoint thật → pred.csv hợp lệ. Pull + kiến trúc + chạy đều đạt.", 12, GREEN, False, True)],
])
page_no(s, 6)

# ============================ 7 · KẾT QUẢ ============================
s = prs.slides.add_slide(BLANK)
header(s, "05 · Kết quả độ chính xác", "Self‑consistency trên mô hình ≤5B đạt 83.59 trên leaderboard; điểm phân bố theo cụm đề.")
bars = [("Toán / Hoá\n(quant)", 73.91, ORANGE), ("Giáo dục CD\n(civics)", 78.67, ORANGE), ("Khoa học\n(Lý…)", 85.41, GREEN)]
base_y = 6.0; scale = 3.3 / 90.0
for i, (lbl, val, col) in enumerate(bars):
    bx = 1.2 + i * 2.6
    bh = val * scale
    box(s, bx, base_y - bh, 1.8, bh, fill=col)
    text(s, bx - 0.3, base_y - bh - 0.5, 2.4, 0.4, [[(f"{val}", 21, col, True, False)]], align=PP_ALIGN.CENTER)
    text(s, bx - 0.3, base_y + 0.1, 2.4, 0.7, [[(lbl, 12, INK, False, False)]], align=PP_ALIGN.CENTER)
box(s, 9.6, 2.4, 3.1, 3.6, fill=INK)
text(s, 9.85, 2.65, 2.65, 3.2, [
    [("83.59", 34, WHITE, True, False)],
    [("leaderboard public‑463", 11.5, CREAM, False, False)],
    [("", 8, INK, False, False)],
    [("80.22", 22, ORANGE, True, False)],
    [("proxy 450 câu (đo nội bộ)", 11, RGBColor(0xC9, 0xCE, 0xD6), False, False)],
])
text(s, 0.7, 6.7, 11.8, 0.4, [[("Lỗi của một mô hình 4B trên đề THPT chủ yếu là sai suy luận/tính toán (cụm quant thấp nhất), không phải thiếu dữ kiện tra‑cứu được — định hướng việc chọn và loại đòn bẩy.", 11.5, GRAY, False, True)]])
page_no(s, 7)

# ============================ 8 · PHƯƠNG PHÁP ============================
s = prs.slides.add_slide(BLANK)
header(s, "06 · Đòn bẩy đang triển khai", "Bốn kỹ thuật ngoại tuyến, một mô hình; mỗi kỹ thuật giữ lại sau khi đo xác nhận hiệu quả.")
levers = [
    ("Chain‑of‑thought tự nhất quán", "Cho mô hình lập luận rồi trích chữ cái cuối; lấy mẫu, bỏ phiếu, và đặt độ tin cậy bằng tỉ lệ đồng thuận. Đây là engine cốt lõi."),
    ("Phán đoán an‑toàn‑từ‑chối", "Nhóm câu 'làm thế nào để [hành vi hại]' có đáp đúng là từ chối; phán đoán theo ngữ nghĩa, không khớp từ khóa nên tổng quát đa ngôn ngữ."),
    ("Constrained decoding (GBNF)", "Ở bước hiệu chỉnh, văn phạm ràng buộc đầu ra về đúng một chữ cái hợp lệ — loại bỏ phỏng đoán heuristic khi một mẫu trôi định dạng."),
    ("Khử thiên lệch vị trí", "Hoán vị tuần hoàn thứ tự lựa chọn trước khi bỏ phiếu, trung hòa thiên lệch vị trí ở các câu nhiều phương án."),
]
for i, (t, d) in enumerate(levers):
    cx = 0.7 + (i % 2) * 6.15
    cy = 2.3 + (i // 2) * 2.15
    ch = 1.9
    box(s, cx, cy, 5.85, ch, fill=WHITE, line=LIGHT, line_w=1.0)
    box(s, cx, cy, 5.85, 0.1, fill=ORANGE)
    text(s, cx + 0.3, cy + 0.12, 5.3, ch - 0.2,
         [[(t, 15, ORANGE_DK, True, False)], [("", 6, INK, False, False)], [(d, 12.5, INK, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space=4.0)
page_no(s, 8)

# ============================ 9 · KỶ LUẬT (điểm nhấn Ý tưởng) ============================
s = prs.slides.add_slide(BLANK)
header(s, "07 · Chống quá khớp", "Phần lớn kỹ thuật 'ngầu' bị loại bằng đo lường — và đó chính là chống overfit cho 2000 câu private.")
table(s, 0.7, 2.05, 12.0, [
    [("Kỹ thuật (đã xây + đo trên Qwen3‑4B)", WHITE, True), ("Kết quả đo", WHITE, True), ("Quyết định / lý do", WHITE, True)],
    ["Fine‑tune v1 (math + legal + MCQ)", ("−4.44", RED, True), "Dữ liệu off‑distribution phá định dạng MCQ (quant 73.9→59.1) — catastrophic forgetting."],
    ["Fine‑tune v2 (chỉ 823 MCQ, hết data)", ("±0.00", GRAY, True), "Base đã ở trần; 823 dòng quá ít để dịch chuyển một mô hình instruct mạnh."],
    ["RAG‑gated (truy hồi tri thức)", ("ẢO / cấm", RED, True), "+3.11 trên proxy là variance k=1; battery sạch civics −5. Luật BTC cấm embedding/rerank."],
    ["TIR — mô hình viết & chạy Python", ("−16.52", RED, True), "4B dưới ngưỡng viết Python đúng (41% câu degrade); phá 24 câu suy luận trực tiếp đã đúng."],
    ["Bỏ phiếu k=5 / tiered", ("≈ 0", GRAY, True), "Trong sai số — self‑consistency đã đủ; thêm mẫu không cải thiện mà tốn thời gian."],
], [4.3, 1.6, 6.1], row_h=0.72, font_sz=11.5, head_sz=12.5)
text(s, 0.7, 6.7, 12, 0.55, [[
    ("Vấn đề thật là variance k=1 (~30 câu lật mỗi lần chạy) → điểm proxy là bẫy; ta đánh giá bằng khả năng tổng quát hóa. ", 12, INK, False, False),
    ("Kết quả phủ định có giá trị ngang kết quả khẳng định — chúng ngăn việc ship phức tạp bất lợi cho tập riêng.", 12, ORANGE_DK, False, True)]])
page_no(s, 9)

# ============================ 10 · THỜI GIAN ============================
s = prs.slides.add_slide(BLANK)
header(s, "08 · Điểm thời gian", "Tối giản chính là nhanh nhất: đường ship một‑mô‑hình vừa thắng độ chính xác vừa thắng thời gian.")
two = [
    ("~3 giây / câu", "≈ 100 phút cho 2000 câu trên một GPU, VRAM ~5GB — không tràn bộ nhớ, không hết giờ trên server 16GB."),
    ("Phức tạp = chậm hơn", "TIR (hai lượt + chạy mã), RAG (truy hồi), k>1 (nhiều mẫu) đều chậm hơn — mà đo ra đều không giúp. Loại chúng bảo vệ cả Time lẫn Accuracy."),
]
for i, (t, d) in enumerate(two):
    cx = 0.7 + i * 6.15
    box(s, cx, 2.35, 5.85, 1.95, fill=CREAM, line=ORANGE, line_w=1.5)
    text(s, cx + 0.32, 2.5, 5.25, 1.7,
         [[(t, 16, ORANGE_DK, True, False)], [("", 6, INK, False, False)], [(d, 12.5, INK, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space=4.0)
text(s, 0.7, 4.7, 12, 1.6, [
    [("Bảo vệ thêm điểm thời gian:", 14, INK, True, False)],
    [("• Phân bổ tính toán theo độ đồng thuận — câu chắc chắn xong trong một lượt rẻ, chỉ phần đuôi bất định mới cần thêm.", 12.5, INK, False, False)],
    [("• Checkpoint mỗi câu + tự khôi phục — container bị ngắt vẫn tiếp tục, không tính toán lại từ đầu.", 12.5, INK, False, False)],
    [("• Một mô hình duy nhất, không khâu truy hồi/giám định ngoài → đường suy luận ngắn và tất định.", 12.5, INK, False, False)],
])
page_no(s, 10)

# ============================ 11 · ĐỊNH VỊ ============================
s = prs.slides.add_slide(BLANK)
header(s, "09 · Định vị kết quả", "83.59 tiệm cận giới hạn thực tế của một mô hình ≤5B; đây là kết luận có dữ liệu, không phải dừng sớm.")
stats = [("4", "đòn bẩy đã đo\n(FT, RAG, TIR, k5)", GRAY), ("0", "đòn bẩy thắng sạch\nngoài self‑consistency", RED), ("100%", "quyết định dựa\nsố liệu đo được", GREEN)]
for i, (n, d, c) in enumerate(stats):
    cx = 0.9 + i * 4.0
    box(s, cx, 2.45, 3.6, 1.7, fill=WHITE, line=c, line_w=1.5)
    text(s, cx, 2.45, 3.6, 1.7,
         [[(n, 27, c, True, False)], [(d, 11.5, INK, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, space=4.0)
text(s, 0.7, 4.55, 12, 1.9, [
    [("Mọi đòn bẩy còn lại đều nằm trong sai số hoặc gây hại khi đo sạch", 14, ORANGE_DK, True, False),
     (" — bằng chứng rằng base Qwen3‑4B đã gần trần của những gì kỹ thuật runtime có thể đạt trên phân phối này.", 13, INK, False, False)],
    [("Đường tăng điểm tiếp theo (nếu có) là thêm dữ liệu trắc nghiệm đa môn — một bài toán thu thập dữ liệu, không phải kỹ thuật runtime.", 13, INK, False, False)],
    [("Tài liệu trình bày số liệu đo được, không cam kết số liệu chưa kiểm chứng.", 13, GREEN, False, True)],
])
page_no(s, 11)

# ============================ 12 · BÀI NỘP / TÁI LẬP ============================
s = prs.slides.add_slide(BLANK)
header(s, "10 · Bài nộp", "Tái lập kết quả trong container bằng hai câu lệnh; mô hình đã đóng gói sẵn bên trong.")
box(s, 0.7, 2.3, 12.0, 1.95, fill=CODEBG)
text(s, 0.95, 2.5, 11.6, 1.75, [
    [("# 1. Tải ảnh tự chứa (mô hình ≤5B đã nướng sẵn bên trong)", 12.5, RGBColor(0x9C, 0xA3, 0xAF), False, False)],
    [("docker pull hacamy12345/neko-core:qwen3-4b-selfconsist-20260618", 13, RGBColor(0x86, 0xEF, 0xAC), False, False)],
    [("# 2. Chạy trên thư mục chứa private_test.csv (hoặc public_test.csv)", 12.5, RGBColor(0x9C, 0xA3, 0xAF), False, False)],
    [("docker run --rm --gpus all -v ./data:/data -v ./output:/output \\", 13, RGBColor(0x86, 0xEF, 0xAC), False, False)],
    [("    hacamy12345/neko-core:qwen3-4b-selfconsist-20260618     (kết quả: ./output/pred.csv)", 13, RGBColor(0x86, 0xEF, 0xAC), False, False)],
])
text(s, 0.7, 4.55, 12, 1.9, [
    [("Mô hình triển khai: Qwen3‑4B‑Instruct‑2507 Q5_K_M (dense 4B, ≤5B tổng), một mô hình, ngoại tuyến, ~5GB VRAM, chạy mọi GPU NVIDIA ≥ Pascal.", 13.5, INK, False, False)],
    [("Quy trình: chain‑of‑thought tự nhất quán + an‑toàn‑từ‑chối + hiệu chỉnh ràng buộc + hiệu chỉnh hợp đồng bất khả xâm phạm.", 13, INK, False, False)],
    [("Ba tag cùng một digest (sha256:39c7891c…): :qwen3-4b-selfconsist-20260618 = :v0.7.2 = :latest. Mã nguồn + tái lập: README.md.", 12, ORANGE_DK, False, True)],
])
page_no(s, 12)

# ============================ 13 · KẾT LUẬN ============================
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, fill=INK)
box(s, 0, 3.05, 13.333, 0.06, fill=ORANGE)
text(s, 0.9, 1.85, 11.5, 1.2, [[("Mô hình suy luận; hệ thống đo lường, bảo vệ khỏi 0 điểm, và chọn kỹ thuật bằng số liệu thực nghiệm.", 25, WHITE, True, False)]])
text(s, 0.9, 3.45, 11.5, 1.7, [
    [("83.59 trên leaderboard với một mô hình ≤5B; bài nộp chạy được trên mọi GPU NVIDIA 2016→2025 (đã verify literal);", 16, CREAM, False, False)],
    [("hợp đồng pred.csv bất khả xâm phạm; và kỷ luật đo lường loại thẳng mọi phức tạp bất lợi cho tập 2000 câu riêng.", 16, CREAM, False, False)],
])
text(s, 0.9, 6.4, 11.5, 0.5, [[("Neko Core — Đại học Hàng hải Việt Nam (VMU) · HackAIthon 2026 · Bảng C (Innovator)", 13, GRAY, False, True)]])

out = Path(__file__).resolve().parents[1] / "docs" / "Neko-Core-Thuyet-minh-phuong-phap.pptx"
try:
    prs.save(str(out))
except PermissionError:
    out = out.with_name("Neko-Core-Thuyet-minh-phuong-phap-new.pptx")
    prs.save(str(out))
    print("(canonical file was locked/open — wrote -new variant)")
print(f"Wrote {out}  ({len(prs.slides._sldIdLst)} slides)")
